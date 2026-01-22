import io
import os
from pptx import Presentation
from pptx.util import Inches, Pt

def fit_image_in_box(slide, img_blob, left, top, max_width, max_height):
    """
    Smartly fits an image into a bounding box while maintaining aspect ratio.
    It prevents distortion and centers the image within the available space.
    """
    if not img_blob: return
    
    stream = io.BytesIO(img_blob)
    
    # 1. Add Picture (constrain by Width first)
    pic = slide.shapes.add_picture(stream, left, top, width=max_width)

    # 2. Check if it's too tall
    if pic.height > max_height:
        # It's too tall! Remove and add by Height instead.
        # We need to access the element to delete it properly
        sp = pic.element
        sp.getparent().remove(sp)

        # Re-add using Height constraint
        stream.seek(0) 
        pic = slide.shapes.add_picture(stream, left, top, height=max_height)

        # Center Horizontally in the box
        if pic.width < max_width:
            center_offset = (max_width - pic.width) / 2
            pic.left = int(left + center_offset)

    else:
        # It fits in height, but might be short. Center Vertically.
        if pic.height < max_height:
            center_offset = (max_height - pic.height) / 2
            pic.top = int(top + center_offset)

def create_presentation(template_path, slides_data):
    """
    Creates a PowerPoint presentation from structured data using a template.
    """
    prs = Presentation(template_path)
    
    # Delete existing slides from the template so we start fresh with the design/master
    # Note: python-pptx doesn't have a simple "delete all slides", we have to remove them via xml
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    for slide in slides:
        xml_slides.remove(slide)

    # Use the first reasonable layout. 
    # Usually index 1 is "Title and Content".
    if len(prs.slide_layouts) > 1:
        slide_layout = prs.slide_layouts[1]
    else:
        slide_layout = prs.slide_layouts[0]

    for i, slide_data in enumerate(slides_data):
        slide = prs.slides.add_slide(slide_layout)
        
        # Data extraction
        title_text = slide_data.get("title", "Untitled")
        content_lines = slide_data.get("content", [])
        
        # Since we don't have image generation in this version yet, disable images
        img_count = 0 
        img_blob = None # Placeholder for future
        
        # 1. Set Title & Calculate Safe Top Margin
        title_shape = slide.shapes.title
        
        # Default safe top (starts at 2.0 inches to be very safe)
        safe_top = Inches(2.0) 

        if title_shape:
            title_shape.text = title_text
            # Calculate where the title visually ends to prevent overlap
            safe_top = title_shape.top + title_shape.height + Inches(0.2)
            # If title is long, add extra buffer
            if len(title_text) > 40: safe_top += Inches(0.4)
            
            # CALCULATE OVERLAP BARRIER
            # We take the title's bottom edge + 0.5 inches buffer
            title_bottom = title_shape.top + title_shape.height
            barrier = title_bottom + Inches(0.5)
            
            # If the calculated barrier is deeper than the default, use the barrier
            if barrier > safe_top:
                safe_top = barrier
                
            # Extra check: If title is very long (wrapping), push it down further
            if len(title_text) > 50:
                 safe_top += Inches(0.5)

        # 2. Add Content (Text Body)
        body_shape = None
        for shape in slide.placeholders:
             if shape.placeholder_format.idx == 1:
                 body_shape = shape
                 break
        
        # Basic Logic (Strategy A/B condensed since we have no images)
        final_body_shape = body_shape
        
        if final_body_shape:
            # CRITICAL OVERLAP FIX
            if final_body_shape.top < safe_top:
                final_body_shape.top = safe_top
            
            # Ensure height doesn't run off slide
            slide_height = prs.slide_height
            remaining_height = slide_height - safe_top - Inches(0.5)
            if final_body_shape.height > remaining_height:
                final_body_shape.height = remaining_height
            
            # Add text
            tf = final_body_shape.text_frame
            tf.clear()
            for point in content_lines:
                p = tf.add_paragraph()
                p.text = point
                p.level = 0
                p.space_before = Pt(6)

                # Readable Font Size
                if p.font.size is None or p.font.size > Pt(32):
                    p.font.size = Pt(20)
                if i == 0: # First slide might be title slide essentially
                    p.font.size = Pt(24)

        # 5. Cleanup Ghost Placeholders (if any interfering)
        for shape in list(slide.placeholders):
            if slide.shapes.title and shape.element == slide.shapes.title.element: continue
            if final_body_shape and shape.element == final_body_shape.element: continue

            try:
                sp = shape.element
                sp.getparent().remove(sp)
            except:
                pass

        # 6. Notes
        if "notes" in slide_data:
            try:
                slide.notes_slide.notes_text_frame.text = slide_data["notes"]
            except:
                pass

    output_filename = "generated_presentation.pptx"
    prs.save(output_filename)
    return output_filename
